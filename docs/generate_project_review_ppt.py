from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        p = text_frame.add_paragraph() if index > 0 else text_frame.paragraphs[0]
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0
        p.text = text
        p.level = level
        p.font.size = Pt(22 if level == 0 else 18)


def add_timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Delivery Timeline (Next 1 Week)"

    left = Inches(0.7)
    top = Inches(1.7)
    width = Inches(12.0)
    height = Inches(0.8)
    bar = slide.shapes.add_shape(1, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(217, 226, 243)
    bar.line.color.rgb = RGBColor(47, 85, 151)

    milestones = [
        ("M1", "Stabilize Dispatch Financial Flow"),
        ("M2", "Harden Migrations and Rollbacks"),
        ("M3", "QA Regression and Release Prep"),
    ]

    start_x = Inches(1.0)
    gap_x = Inches(3.9)
    for i, (code, label) in enumerate(milestones):
        x = start_x + Inches(i * gap_x)
        shape = slide.shapes.add_textbox(x, Inches(3.0), Inches(3.8), Inches(1.8))
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = code
        p1.font.bold = True
        p1.font.size = Pt(24)
        p1.alignment = PP_ALIGN.LEFT

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(16)
        p2.level = 1


def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Portal ERP - Project Review",
        "Date: 2026-05-06 | Focus: Dispatch, Payments, Batch Workflow",
    )

    add_bullet_slide(
        prs,
        "Executive Snapshot",
        [
            "Sprint theme: Dispatch financial workflow and operational consolidation",
            "Backend readiness: High | Frontend readiness: Medium-High",
            "QA readiness: Medium | Release readiness: Not yet",
            "Overall status: In progress with strong momentum",
        ],
    )

    add_bullet_slide(
        prs,
        "Scope Delivered in Current Cycle",
        [
            "Dispatch flow",
            ("FY-based dispatch numbering and transactional create/update flow", 1),
            ("Nested financial and status payload persistence", 1),
            ("Integrated immediate/partial payment capture", 1),
            "Data model and migrations",
            ("New dispatch payments table and Eloquent model", 1),
            ("Payment methods table upgraded with metadata and audit fields", 1),
            "Batch operations",
            ("Stock adjustment and work-order production refresh improvements", 1),
        ],
    )

    add_bullet_slide(
        prs,
        "Codebase Change Highlights",
        [
            "Backend: DispatchController, BatchController, DispatchStoreRequest, Dispatch model",
            "New models: DispatchFinancial, DispatchPayment",
            "New migrations: mm_dispatch_payments + payment methods enhancement",
            "Frontend: major updates in Batches/Dispatches components and supporting forms",
            "Build output regenerated (large diff under public/build)",
        ],
    )

    add_bullet_slide(
        prs,
        "Risks and Controls",
        [
            "Migration rollback safety for altered payment method table",
            "Payload contract drift between Vue forms and backend request mapping",
            "Payment update semantics for multi-payment scenarios per dispatch",
            "PR/review noise from generated assets mixed with source changes",
            "Mitigation: tests, schema freeze, rollback hardening, PR split strategy",
        ],
    )

    add_timeline_slide(prs)

    add_bullet_slide(
        prs,
        "Next 3 Execution Steps",
        [
            "Freeze frontend/backend payload schema for dispatch forms",
            "Add feature tests for dispatch create/update and payment capture",
            "Validate migration down-paths and FK-safe rollback behavior",
            "Split source changes and build artifacts into separate review units",
        ],
    )

    add_bullet_slide(
        prs,
        "Definition of Done",
        [
            "Dispatch lifecycle works end-to-end with financials, status, and payments",
            "Migrations pass reliable up/down verification in staging",
            "Batch-to-work-order stock and production numbers are regression-tested",
            "QA sign-off completed for happy-path and failure-path scenarios",
            "Release package is clean, reviewable, and deployment-ready",
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("docs/project-review-2026-05-06.pptx")
