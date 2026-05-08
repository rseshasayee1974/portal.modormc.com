from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def set_slide_title(slide, title_text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 56, 100)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_title(slide, "Portal ERP - Project Flow Overview")

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(2.5))
    tf = sub.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "Client Review Presentation"
    p1.font.size = Pt(34)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(47, 85, 151)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Date: 06 May 2026"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(68, 84, 106)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Objective: Show end-to-end workflow, delivery status, and next milestones"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(68, 84, 106)
    p3.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_title(slide, title)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(5.8))
    tf = body.text_frame
    tf.clear()

    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            txt, level = item
        else:
            txt, level = item, 0
        p.text = txt
        p.level = level
        p.font.size = Pt(24 if level == 0 else 18)
        p.font.color.rgb = RGBColor(34, 34, 34)


def add_flow_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_title(slide, "End-to-End Project Flow")

    steps = [
        "1. Enquiry\n& Planning",
        "2. Work Order\nCreation",
        "3. Batch\nProduction",
        "4. Dispatch\nExecution",
        "5. Delivery\nConfirmation",
        "6. Billing\n& Collection",
        "7. Dashboard\n& Reporting",
    ]

    x_start = 0.45
    y = 2.3
    w = 1.65
    h = 1.1
    gap = 0.22

    for i, step in enumerate(steps):
        x = Inches(x_start + i * (w + gap))
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(y), Inches(w), Inches(h)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(221, 235, 247)
        rect.line.color.rgb = RGBColor(47, 85, 151)

        tf = rect.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 56, 100)
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(
                Inches(x_start + i * (w + gap) + w + 0.03), Inches(y + 0.35), Inches(0.15), Inches(0.3)
            )
            ap = arrow.text_frame.paragraphs[0]
            ap.text = ">"
            ap.font.size = Pt(20)
            ap.font.bold = True
            ap.font.color.rgb = RGBColor(91, 155, 213)
            ap.alignment = PP_ALIGN.CENTER

    note = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.6), Inches(1.8))
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.text = "Current sprint focus is Step 4 to Step 6 (Dispatch, Payments, and Billing readiness)."
    np.font.size = Pt(18)
    np.font.color.rgb = RGBColor(68, 84, 106)
    np.alignment = PP_ALIGN.CENTER


def add_status_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_title(slide, "Current Delivery Status")

    columns = [
        ("Completed", RGBColor(198, 239, 206), [
            "Dispatch numbering logic",
            "Dispatch status capture",
            "Payment model and migration",
            "Batch workflow integration",
        ]),
        ("In Progress", RGBColor(255, 235, 156), [
            "Frontend dispatch form alignment",
            "Validation contract hardening",
            "End-to-end QA scenarios",
        ]),
        ("Next", RGBColor(221, 235, 247), [
            "Migration rollback hardening",
            "Client UAT feedback loop",
            "Go-live checklist",
        ]),
    ]

    x_positions = [0.6, 4.75, 8.9]
    for idx, (title, color, items) in enumerate(columns):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x_positions[idx]),
            Inches(1.4),
            Inches(3.7),
            Inches(4.9),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = RGBColor(89, 89, 89)

        tf = card.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(21)
        p0.font.bold = True
        p0.alignment = PP_ALIGN.CENTER

        for it in items:
            p = tf.add_paragraph()
            p.text = f"- {it}"
            p.font.size = Pt(15)
            p.level = 1


def add_timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_title(slide, "Client Timeline and Milestones")

    milestones = [
        ("Week 1", "Stabilize dispatch + payment flow"),
        ("Week 2", "Regression testing and UAT support"),
        ("Week 3", "Release candidate and deployment prep"),
    ]

    for i, (wk, task) in enumerate(milestones):
        y = 1.5 + (i * 1.8)
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.0), Inches(y), Inches(1.2), Inches(1.2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(47, 85, 151)
        circle.line.color.rgb = RGBColor(47, 85, 151)
        ctf = circle.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.text = wk
        cp.font.size = Pt(13)
        cp.font.bold = True
        cp.font.color.rgb = RGBColor(255, 255, 255)
        cp.alignment = PP_ALIGN.CENTER

        task_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.6), Inches(y), Inches(9.8), Inches(1.2)
        )
        task_box.fill.solid()
        task_box.fill.fore_color.rgb = RGBColor(242, 242, 242)
        task_box.line.color.rgb = RGBColor(200, 200, 200)
        ttf = task_box.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        tp.text = task
        tp.font.size = Pt(20)
        tp.font.color.rgb = RGBColor(31, 56, 100)


def add_closing_slide(prs):
    add_bullet_slide(
        prs,
        "Client Review - Required Inputs",
        [
            "Confirm final approval flow for dispatch to billing",
            "Validate payment capture options (full / partial)",
            "Share UAT test users and plant-wise review schedule",
            "Approve target go-live week after UAT sign-off",
        ],
    )


def build_ppt(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "Project Objective and Scope",
        [
            "Objective: Streamline order-to-cash flow for RMC operations",
            "Scope in this phase:",
            ("Work order, batching, dispatch, payment, and billing handoff", 1),
            ("Operational visibility through dashboard and status tracking", 1),
            ("Data consistency across plant operations", 1),
        ],
    )
    add_flow_overview_slide(prs)
    add_status_slide(prs)
    add_timeline_slide(prs)
    add_closing_slide(prs)

    prs.save(output_path)


if __name__ == "__main__":
    build_ppt("docs/project-flow-overview-client-2026-05-06.pptx")
